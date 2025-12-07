"""
File Processing Routes - File upload and processing functionality
"""
from flask import Blueprint, request, jsonify
from datetime import datetime
from limiter_config import get_limiter
import tempfile
import os
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pdfplumber
import re
import io

# Create blueprint
file_bp = Blueprint('file', __name__)

# Get limiter instance
limiter = get_limiter(None)  # Will be configured in main app

def process_pdf_file(file_data, filename):
    """
    Process PDF file using pdfplumber and extract filtered API-focused text with highlights
    """
    # API-related keywords to search for (case-insensitive)
    api_keywords = [
        "authentication", "authorization", "token", "API key", "endpoint", 
        "base URL", "GET", "POST", "PUT", "DELETE", "parameters", "response", 
        "status code", "example", "error", "request", "header", "body"
    ]
    
    # Extended API-related patterns for filtering
    api_patterns = [
        r'\b(GET|POST|PUT|DELETE|PATCH)\b',  # HTTP methods
        r'curl',  # curl commands
        r'api|url|http|https',  # API/URL terms
        r'endpoint|parameter|request|response',  # API concepts
        r'status|code|error|success|fail',  # Status codes
        r'json|xml|format|type',  # Data formats
        r'header|body|content-type',  # Request/response parts
        r'example|usage|documentation',  # Documentation terms
        r'[{}()[\]<>]',  # Code brackets
        r'["\']\w+["\']\s*:',  # JSON-like structures
    ]
    
    try:
        
        # Get file extension
        file_extension = filename.lower().split('.')[-1] if '.' in filename else ''
        
        # Create a temporary file with appropriate extension
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            temp_file.write(file_data)
            temp_file_path = temp_file.name
        
        
        # Try to extract text from the file
        filtered_text = ""
        highlighted_lines = []
        total_lines = 0
        kept_lines = 0
        
        try:
            with pdfplumber.open(temp_file_path) as pdf:
                
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        # Split page text into lines and filter
                        lines = page_text.split('\n')
                        page_filtered_lines = []
                        
                        for line_num, line in enumerate(lines):
                            line = line.strip()
                            total_lines += 1
                            
                            # Skip empty lines
                            if not line:
                                continue
                            
                            # Skip very short lines (less than 10 characters)
                            if len(line) < 10:
                                continue
                            
                            # Check if line contains API-related content
                            line_lower = line.lower()
                            is_api_related = False
                            
                            # Check API keywords
                            for keyword in api_keywords:
                                if keyword.lower() in line_lower:
                                    is_api_related = True
                                    break
                            
                            # Check API patterns
                            if not is_api_related:
                                for pattern in api_patterns:
                                    if re.search(pattern, line, re.IGNORECASE):
                                        is_api_related = True
                                        break
                            
                            # If line is API-related, keep it
                            if is_api_related:
                                page_filtered_lines.append(line)
                                kept_lines += 1
                                
                                # Add to highlights if it contains API keywords
                                for keyword in api_keywords:
                                    if keyword.lower() in line_lower:
                                        highlighted_line = f"Page {page_num + 1}, Line {line_num + 1}: {line}"
                                        if highlighted_line not in highlighted_lines:
                                            highlighted_lines.append(highlighted_line)
                                        break
                            else:
                        
                        # Add filtered lines to extracted text
                        if page_filtered_lines:
                            page_filtered_text = '\n'.join(page_filtered_lines)
                            filtered_text += page_filtered_text + "\n"
                        else:
                    else:
        except Exception as pdf_error:
            if file_extension != 'pdf':
                return {"error": f"This file type ({file_extension.upper()}) cannot be processed. Please upload a PDF file."}
            else:
                return {"error": f"PDF processing failed: {str(pdf_error)}. Please try with a different file."}
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        if filtered_text.strip():
            
            return {
                "text": filtered_text.strip(), 
                "success": True,
                "highlights": highlighted_lines
            }
        else:
            return {"error": "No API-related content could be extracted from the file. Please try with a different file."}
            
    except Exception as e:
        # Clean up temporary file if it exists
        try:
            if 'temp_file_path' in locals():
                os.unlink(temp_file_path)
        except:
            pass
        return {"error": f"File processing error: {str(e)}. Please try again later."}

def process_docx_file(file_data, filename):
    """Process Word documents (.docx) and extract text"""
    try:
        
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(file_data)
            temp_file_path = temp_file.name
        
        # Read the document
        doc = Document(temp_file_path)
        
        # Extract text from all paragraphs
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        full_text = '\n'.join(text_content)
        
        return {
            "text": full_text,
            "success": True
        }
        
    except Exception as e:
        return {"error": f"Failed to process Word document: {str(e)}"}

def process_excel_file(file_data, filename):
    """Process Excel files (.xlsx, .xls) and extract text"""
    try:
        
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
            temp_file.write(file_data)
            temp_file_path = temp_file.name
        
        # Read the Excel file
        workbook = load_workbook(temp_file_path, data_only=True)
        
        text_content = []
        
        # Process each sheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content.append(f"\n--- Sheet: {sheet_name} ---\n")
            
            # Get all values from the sheet
            for row in sheet.iter_rows(values_only=True):
                row_data = []
                for cell_value in row:
                    if cell_value is not None:
                        row_data.append(str(cell_value))
                if row_data:
                    text_content.append(' | '.join(row_data))
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        full_text = '\n'.join(text_content)
        
        return {
            "text": full_text,
            "success": True
        }
        
    except Exception as e:
        return {"error": f"Failed to process Excel file: {str(e)}"}

def process_pptx_file(file_data, filename):
    """Process PowerPoint files (.pptx) and extract text"""
    try:
        
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(file_data)
            temp_file_path = temp_file.name
        
        # Read the presentation
        presentation = Presentation(temp_file_path)
        
        text_content = []
        
        # Process each slide
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"\n--- Slide {slide_num} ---\n")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        full_text = '\n'.join(text_content)
        
        return {
            "text": full_text,
            "success": True
        }
        
    except Exception as e:
        return {"error": f"Failed to process PowerPoint file: {str(e)}"}

def process_csv_file(file_data, filename):
    """Process CSV files and extract text"""
    try:
        
        # Decode the file data
        text_content = file_data.decode('utf-8')
        
        # Try to parse as CSV and convert to readable format
        try:
            df = pd.read_csv(io.StringIO(text_content))
            # Convert DataFrame to string representation
            csv_text = df.to_string(index=False)
            return {
                "text": csv_text,
                "success": True
            }
        except Exception as csv_error:
            # If pandas fails, return the raw text
            return {
                "text": text_content,
                "success": True
            }
        
    except Exception as e:
        return {"error": f"Failed to process CSV file: {str(e)}"}

def process_text_file(file_data, filename):
    """Process text files (.txt, .md, .json, .xml, .rtf) and extract text"""
    try:
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                text_content = file_data.decode(encoding)
                return {
                    "text": text_content,
                    "success": True
                }
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, try with error handling
        text_content = file_data.decode('utf-8', errors='ignore')
        return {
            "text": text_content,
            "success": True
        }
        
    except Exception as e:
        return {"error": f"Failed to process text file: {str(e)}"}

@file_bp.route('/file-to-text', methods=['POST', 'OPTIONS'])
def file_to_text_endpoint():
    # Handle CORS preflight request
    if request.method == 'OPTIONS':
        return '', 204
    
    # For POST requests, call the rate-limited function
    return file_to_text_post()

@limiter.limit("10 per minute")
def file_to_text_post():
    """
    Process uploaded files and extract text (PDF, images, etc.)
    """
    try:
        
        # Check if file was uploaded
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        
        # Check if file is empty
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Read file data
        file_data = file.read()
        
        # Determine file type and process accordingly
        file_extension = file.filename.lower().split('.')[-1] if '.' in file.filename else ''
        
        if file_extension == 'pdf':
            # Process PDF with pdfplumber
            result = process_pdf_file(file_data, file.filename)
        elif file_extension in ['docx', 'doc']:
            # Process Word documents
            result = process_docx_file(file_data, file.filename)
        elif file_extension in ['xlsx', 'xls']:
            # Process Excel files
            result = process_excel_file(file_data, file.filename)
        elif file_extension in ['pptx', 'ppt']:
            # Process PowerPoint files
            result = process_pptx_file(file_data, file.filename)
        elif file_extension == 'csv':
            # Process CSV files
            result = process_csv_file(file_data, file.filename)
        elif file_extension in ['txt', 'md', 'json', 'xml', 'rtf']:
            # Process text files
            result = process_text_file(file_data, file.filename)
        else:
            return jsonify({'error': f'Unsupported file type: {file_extension}. Supported formats: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), CSV, and text files (.txt, .md, .json, .xml, .rtf).'}), 400
        
        
        if result.get('success'):
            response_data = {
                'success': True,
                'text': result['text'],
                'message': f'Text extracted successfully from {file_extension.upper()} file'
            }
            
            # Add highlights if available
            if 'highlights' in result:
                response_data['highlights'] = result['highlights']
            
            return jsonify(response_data)
        else:
            return jsonify({
                'success': False,
                'error': result.get('error', 'Unknown error occurred')
            }), 400
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'Server error: {str(e)}'
        }), 500